import datetime
import matplotlib.pyplot as plt
from numpy import negative
import squarify
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
import copy, io
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, Field
from typing import List, Optional
import pandas as pd

class Filter(BaseModel):
    field: str
    value: str

class Timeframe(BaseModel):
    from_: int = Field(alias="from")
    to: int 
    timeframe_field: Optional[str]

class ReportParams(BaseModel):
    id: str
    reportType: str
    filters: List[Filter] 
    timeframe: Timeframe 

def parse_report_params(data: dict) -> ReportParams:
    return ReportParams(
        id=data.get("id"),
        reportType=data.get("reportType"),
        filters=[Filter(**f) for f in data.get("filters", [])],
        timeframe=Timeframe(**data.get("timeframe")),
    )


def replace_chart(slide, C_index, C_data):
    chart = 0
    for shape in slide.shapes:
        if shape.has_chart:
            chart += 1
            if chart == C_index + 1:
                chart = shape.chart
                chart.replace_data(C_data)
                print(f"Chart index {C_index} replaced")
                return
    print(f"Chart index {C_index} not found")

def generate_time():
    current_time = datetime.datetime.now()
    epoch_millis = int((current_time - datetime.datetime(1970, 1, 1)).total_seconds() * 1000)

    return epoch_millis

def create_treemap(data, treemap_image_path='treemap_custom.png'):
    # Values and labels
    values = []
    labels = []
    for item in data:
        if item['value'] is not None and item['value'] > 0:
            values.append(item['value'])
            labels.append(item['text'])
    
    if not values:  # If there are no valid values
        print("No valid values to plot.")
        return
    
    # Normalize values to ensure they fit within the plotting area
    total_value = sum(values)
    values = [value / total_value for value in values]
    
    # Formatted labels
    formatted_labels = [f"{label}, {value*100:.2f}%" for label, value in zip(labels, values)]
    
    # Colors for each category
    colors = ['#bdbdbd', '#e41a1c', '#377eb8'] * (len(values) // 3 + 1)  # Ensure enough colors
    
    # Calculate the positions for the treemap
    rects = squarify.squarify(sizes=values, x=0, y=0, dx=1, dy=1)
    
    # Create the plot
    fig, ax = plt.subplots(figsize=(6, 6))
    
    # Plot each rectangle
    for rect, label, color in zip(rects, formatted_labels, colors):
        x, y, w, h = rect['x'], rect['y'], rect['dx'], rect['dy']
        ax.add_patch(plt.Rectangle((x, y), w, h, color=color, ec="white"))
        ax.text(x + 0.01, y + 0.01, label, ha='left', va='bottom', color='white', fontsize=13, weight='bold')
    
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    
    # Save the plot to a file
    plt.savefig(treemap_image_path, bbox_inches='tight', transparent=True, pad_inches=0)
    plt.close()


def MoveSlide(presentation, old_index, new_index):
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def DuplicateSlide(presentation, slide_index):

    # memastikan indeks slide berada dalam rentang yang benar
    if slide_index < 0 | slide_index >= len(presentation.slides):
        raise IndexError("Indeks slide berada di luar rentang yang valid.")

    # ambil slide berdasarkan indeks
    slide = presentation.slides[slide_index]
    slide_layout = slide.slide_layout

    # tambah slide baru dengan layout yang sama
    new_slide = presentation.slides.add_slide(slide_layout)

    # copy semua shape dari slide lama ke slide baru
    for shape in slide.shapes:
        if shape.shape_type == 13:  # jika shape adalah image, jadi 13 adalah kode untuk image 
            # ambil image dari slide lama
            image_stream = io.BytesIO(shape.image.blob)
            # menambahkan image ke slide baru
            new_slide.shapes.add_picture(
                image_stream, shape.left, shape.top, shape.width, shape.height
            )
        else:
            # jika shape non-image, dilakukan deep copy seperti sebelumnya
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # mengosongkan teks pada shape judul agar tampilan menjadi clear
    if new_slide.shapes.title:
        new_slide.shapes.title.text = " "

    return new_slide

def Delete_slide(prs, slide_index):
    """Delete a slide from a PowerPoint presentation by its index."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError("Slide index out of range.")
    slide_id = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[slide_index]
    
def Safe_access(list_data, index):
        return list_data[index] if list_data and len(list_data) > index else {}

# prs = Presentation('[Pemkot Yogyakarta] Report ISA Periode 02-08 Juni 2024.pptx')

# try:
#     duplicate_slide(prs, 5)

#     move_slide(prs, 29, 6)

#     prs.save('contoh_duplicated.pptx')

#     print("Slide berhasil diduplikasi dan disimpan.")
# except IndexError as e:
#     print(f"Terjadi kesalahan: {e}")
# except Exception as e:
#     print(f"Terjadi kesalahan tak terduga: {e}")


def chart_bar_hori(prs, data, slide_index=0, chart_index=0):
    if not data:
        categories = []
        values = []
    else:
        items = [(list(item.keys())[0], list(item.values())[0]) for item in data]
        sorted_data = sorted(items, key=lambda x: x[1], reverse=False)
        categories = [item[0] for item in sorted_data]
        values = [item[1] for item in sorted_data]

    chart_data = CategoryChartData()
    if categories:
        chart_data.categories = categories
        chart_data.add_series('', values)
    else:
        chart_data.categories = ['']
        chart_data.add_series('', [0])

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart(slide, chart_index, chart_data)
            break


def chart_bar_verSent(presentation, data, slide_index=0, chart_index=0):

    categories = []
    positive_values = []
    neutral_values = []
    negative_values = []

    chart_data = CategoryChartData()
    if not data or all(
        int(item.get('positive', 0) or 0) == 0 and 
        int(item.get('netral', 0) or 0) == 0 and 
        int(item.get('negative', 0) or 0) == 0 
        for item in data):
        
        chart_data.categories = ['']
        chart_data.add_series('Positive', [0])
        chart_data.add_series('Neutral', [0])
        chart_data.add_series('Negative', [0])
    else:
        for item in data:
            categories.append(item["media"])
            
            positive_value = int(item.get('positive', 0) or 0)
            neutral_value = int(item.get('netral', 0) or 0)
            negative_value = int(item.get('negative', 0) or 0)

            positive_values.append(positive_value if positive_value != 0 else None)
            neutral_values.append(neutral_value if neutral_value != 0 else None)
            negative_values.append(negative_value if negative_value != 0 else None)

        chart_data.categories = categories
        chart_data.add_series('Positive', positive_values)
        chart_data.add_series('Neutral', neutral_values)
        chart_data.add_series('Negative', negative_values)


    for i, slide in enumerate(presentation.slides):
        if i == slide_index:
            replace_chart(slide, chart_index, chart_data)
            break


def chart_pie_B(presentation, data, slide_index=0, chart_index=0):
    categories = []
    values = []

    if not data:
        total_positive = 0
        categories.append(' ')
        values.append(total_positive)
    else:
        total_positive = sum(item.get('positive', 0) for item in data)
        total_neutral = sum(item.get('neutral', 0) for item in data)
        total_negative = sum(item.get('negative', 0) for item in data)

        if total_positive != 0:
            categories.append('Positive')
            values.append(total_positive)
        
        if total_neutral != 0:
            categories.append('Neutral')
            values.append(total_neutral)
        
        if total_negative != 0:
            categories.append('Negative')
            values.append(total_negative)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)

    for i, slide in enumerate(presentation.slides):
        if i == slide_index:
            replace_chart(slide, chart_index, chart_data)
            break


def chart_pie_K(presentation, data, slide_index=0, chart_index=0, categories='A', value='Positive', non_value='Non Positive'):
    chart_data = CategoryChartData()

    positive_values = []
    neutral_values = []

    if not data or all(
        int(item.get(value, 0) or 0) == 0 and 
        int(item.get(non_value, 0) or 0) == 0
        for item in data):
        
        chart_data.categories = [categories]
        chart_data.add_series(value, [0])
        chart_data.add_series(non_value, [0])
    else:
        for item in data:
            positive_value = int(item.get(value, 0) or 0)
            neutral_value = int(item.get(non_value, 0) or 0)

            positive_values.append(positive_value if positive_value != 0 else None)
            neutral_values.append(neutral_value if neutral_value != 0 else None)

        chart_data.categories = [categories]
        chart_data.add_series('A', [x for pair in zip(positive_values, neutral_values) for x in pair])

    for i, slide in enumerate(presentation.slides):
        if i == slide_index:
            replace_chart(slide, chart_index, chart_data)
            break


def format_date(date_time):

    bulan = {
    "January": "Januari",
    "February": "Februari",
    "March": "Maret",
    "April": "April",
    "May": "Mei",
    "June": "Juni",
    "July": "Juli",
    "August": "Agustus",
    "September": "September",
    "October": "Oktober",
    "November": "November",
    "December": "Desember"
    }
    # Parsing tanggal dari string
    date_from = datetime.strptime(date_time["from"], "%Y-%m-%d %H:%M:%S")
    date_to = datetime.strptime(date_time["to"], "%Y-%m-%d %H:%M:%S")

    # Konversi ke format yang diinginkan
    formatted_date_from = date_from.strftime("%d %B %Y")
    formatted_date_to = date_to.strftime("%d %B %Y")
    time_from = date_from.strftime("%H.%M.%S")
    time_to = date_to.strftime("%H.%M.%S")
    
    # Ganti nama bulan ke bahasa Indonesia
    for english, indonesian in bulan.items():
        formatted_date_from = formatted_date_from.replace(english, indonesian)
        formatted_date_to = formatted_date_to.replace(english, indonesian)
    
    # Membuat dictionary hasil
    formatted_result = {
        "date_from": formatted_date_from,
        "date_to": formatted_date_to,
        "time_from": time_from,
        "time_to": time_to
    }
    
    return formatted_result


def remove_newlines(data):
    if isinstance(data, str):
        return data.replace("\n", "")
    elif isinstance(data, list):
        return [remove_newlines(item) for item in data]
    elif isinstance(data, dict):
        return {key: remove_newlines(value) for key, value in data.items()}
    return data


def retrieve_data(data, key, sub_key=None, default_fill='', length_min=10):
    """
    Contoh penggunaan:

    texts = retrieve_values(data, 'text')

    locations = retrieve_values(data, 'location')

    opds = retrieve_values(data, 'opd', length_min=10)

    values = retrieve_values(data, 'value', default_fill='0')

    sentiments = retrieve_values(data, 'sentiment', sub_key=None, length_min=2)

    insights = retrieve_values(data, 'media', sub_key='text', length_min=5)

    ranks = retrieve_values(data, 'rank', default_fill='0', length_min=3)

    """
    data = remove_newlines(data)

    values = []

    # Ambil item berdasarkan key utama
    items = data.get(key, [])

    if isinstance(items, list):
        for item in items:
            if isinstance(item, dict):
                if sub_key:  # Ambil berdasarkan sub_key jika ada
                    value = item.get(sub_key, default_fill)
                else:  # Ambil item langsung jika tidak ada sub_key
                    value = item.get(key, default_fill)
                values.append(value if value else default_fill)
            elif isinstance(item, (str, int)):
                values.append(str(item))
            elif isinstance(item, list):
                values.extend(item if item else [default_fill])
            else:
                values.append(default_fill)
    elif isinstance(items, dict):
        if sub_key:  # Ambil berdasarkan sub_key jika ada
            value = items.get(sub_key, default_fill)
        else:  # Ambil item langsung jika tidak ada sub_key
            value = items.get(key, default_fill)
        values.append(value if value else default_fill)
    else:
        values = [default_fill]

    # Meratakan list jika ada nested list
    flat_values = []
    for val in values:
        if isinstance(val, list):
            flat_values.extend(val)
        else:
            flat_values.append(val)

    # Pastikan panjang list sesuai dengan min_length
    while len(flat_values) < length_min:
        flat_values.append(default_fill)

    return flat_values


def retrieve_rank(data, key):
    values = []

    value = data.get(key, '')

    if isinstance(value, list):
        for item in value:
            values.append(item if item else '0')
    elif isinstance(value, int):
        values.append(str(value))
    elif isinstance(value, str):
        values.append(value)
    else:
        values = ['0']

    while len(values) < 1:
        values.append('0')

    return values

def ConvertDate(date):
    if "C" in date:
        date = date.replace("C ", "")
    
    tanggal_obj = datetime.strptime(date, "%Y-%m-%d")
    tanggal_ubah = tanggal_obj.strftime("%d %B %Y")
    
    bulan_dict = {
        "January": "Januari", "February": "Februari", "March": "Maret", "April": "April",
        "May": "Mei", "June": "Juni", "July": "Juli", "August": "Agustus",
        "September": "September", "October": "Oktober", "November": "November", "December": "Desember"
    }

    for english, indonesian in bulan_dict.items():
        tanggal_ubah = tanggal_ubah.replace(english, indonesian)
    
    return tanggal_ubah

def detect_chart_indexes(presentation_path):
    prs = Presentation(presentation_path)
    for slide_index, slide in enumerate(prs.slides):
        print(f"Slide {slide_index + 1}")
        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_chart:
                print(f"  Chart found at shape index {shape_index}")

def extract_aggregation_data(aggregation_result):
    try:
        # Safely access the buckets in the aggregation result
        buckets = aggregation_result.get("aggregations", {}).get("2", {}).get("buckets", [])
        if not isinstance(buckets, list):
            raise ValueError("Buckets should be a list.")

        # Extract data from each bucket
        extracted_data = []
        for bucket in buckets:
            extracted_data.append({
                "key": bucket.get("key", ""),                 # Default to "N/A" if key is missing
                "doc_count": bucket.get("doc_count", 0),         # Default to 0 if doc_count is missing
                "value": bucket.get("1", {}).get("value", None)  # Default to None if value is missing
            })
        return extracted_data

    except Exception as e:
        # Log the error and return an empty list
        print(f"Error extracting aggregation data: {e}")
        return []

def retrieve_data_object(data, key=None, sub_key=None, default_fill='', length_min=10):
    """
    Retrieve values from a list of dictionaries or a nested data structure.

    Parameters:
    - data: A list of dictionaries or a dictionary.
    - key: Key to access the list of dictionaries (optional for list input).
    - sub_key: Key within each dictionary to extract values (optional).
    - default_fill: Default value to use for missing data.
    - length_min: Minimum length of the returned list.

    Returns:
    - A list of extracted values, padded to length_min if necessary.
    """
    # If input is a dictionary, extract the list using the 'key' argument
    if isinstance(data, dict) and key:
        items = data.get(key, [])
    elif isinstance(data, list):  # If input is already a list, use it directly
        items = data
    else:
        return [default_fill] * length_min

    values = []

    # Extract values from each item in the list
    if isinstance(items, list):
        for item in items:
            if isinstance(item, dict):
                if sub_key:  # Extract based on sub_key if provided
                    value = item.get(sub_key, default_fill)
                else:  # Default to extracting the whole dictionary if no sub_key
                    value = item
                values.append(value if value else default_fill)
            elif isinstance(item, (str, int)):
                values.append(str(item))
            elif isinstance(item, list):
                values.extend(item if item else [default_fill])
            else:
                values.append(default_fill)
    else:
        values = [default_fill]

    # Flatten the list if nested lists exist
    flat_values = []
    for val in values:
        if isinstance(val, list):
            flat_values.extend(val)
        else:
            flat_values.append(val)

    # Ensure the list is at least length_min
    while len(flat_values) < length_min:
        flat_values.append(default_fill)

    return flat_values

def handle_unested(result):
    buckets = result["aggregations"]["2"]["buckets"]
    return [
        {
            "key": b["key"],
            "doc_count": b["doc_count"]
        }
        for b in buckets
    ]

def handle_unested_v2(result):
    buckets = result["aggregations"]["1"]["buckets"]
    return [
        {
            "key": b["key"],
            "doc_count": b["doc_count"]
        }
        for b in buckets
    ]

def handle_nested(result):
    buckets = result["aggregations"]["2"]["buckets"]
    return [
        {
            "key_as_string": b["key_as_string"],
            "key": b["key"],
            "doc_count": b["doc_count"],
            "inner_buckets": [
                {
                    "key": inner_b["key"],
                    "doc_count": inner_b["doc_count"]
                }
                for inner_b in b["3"]["buckets"]
            ] if "3" in b else []
        }
        for b in buckets
    ]
# def handle_nested(result):
#     buckets = result["aggregations"]["2"]["buckets"]
#     return [
#         {
#             "key_as_string": b["key_as_string"],
#             "key": b["key"],
#             "doc_count": b["doc_count"],
#             "inner_buckets": {
#                 inner_b["key"]: {
#                     "doc_count": inner_b["doc_count"]
#                 }
#                 for inner_b in b["3"]["buckets"]
#             } if "3" in b else {}
#         }
#         for b in buckets
#     ]


def handle_aggregation(result): 
    buckets = result["aggregations"]["2"]["buckets"]
    return [
        {
            "key": b["key"],
            "doc_count": b["doc_count"],
            "value": b["1"]["value"] if "1" in b else None 
        }
        for b in buckets
    ]

def handle_aggregation_v2(result):
    buckets = result["aggregations"]["1"]["buckets"]
    return [
        {
            "key": b["key"],
            # "key_as_string": b["key_as_string"],
            "doc_count": b["doc_count"],
        }
        for b in buckets
    ]

def handle_simple_data(result):
    aggregation = result.get("aggregations")
    return {
        "engagement": aggregation.get("Engagement", {}).get("value", 0),
        "eksposure": aggregation.get("Eksposure", {}).get("value", 0)
    }

def get_nested(datas, length_min):
    keys=[]
    doc_count=[]
    for item in datas:
        inner_buckets=item.get("inner_buckets", {})
        for platform, metrics in inner_buckets.items():
            keys.append(platform)
            doc_count.append(metrics.get("doc_count", 0))
    while len(keys) < length_min:
        keys.append("")
        doc_count.append(0)
    return keys, doc_count

def generate_line_chart(presentation, data, s_idx, c_idx):
    data_chart = data
    
    if not data_chart:
        categories = []
        series_data = {}
    else:
        categories = []
        series_data = {}

    platforms = set()

    # First, collect all platform names
    for item in data_chart:
        for platform in item.get("inner_buckets", []):
            platforms.add(platform["key"])

    platforms = list(platforms)

    for item in data_chart:
        # Add the date as a category
        categories.append(item["key_as_string"])

        # Initialize all platforms for the current date with 0
        daily_data = {platform: 0 for platform in platforms}

        # Update the platforms with actual values
        for platform in item.get("inner_buckets", []):
            daily_data[platform["key"]] = platform["doc_count"]

        # Append the values to the respective platform in series_data
        for platform in platforms:
            if platform not in series_data:
                series_data[platform] = []
            series_data[platform].append(daily_data[platform])

    # Generate the chart data
    chart_data = CategoryChartData()
    if categories:
        chart_data.categories = categories
        for platform, values in series_data.items():
            chart_data.add_series(platform, values)
    else:
        chart_data.categories = [""]
        chart_data.add_series("", [0])

    # Replace the chart in the presentation
    slide_index = s_idx
    chart_index_to_replace = c_idx

    for i, slide in enumerate(presentation.slides):
        if i == slide_index:
            replace_chart(slide, chart_index_to_replace, chart_data)
            break


def find_chart(presentation):
    for slide_index, slide in enumerate(presentation.slides):
        print(f"Slide {slide_index + 1}:")
        
        chart_found = False
        
        chart_index = 0
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                print(f"  - Chart pada indeks chart {chart_index}")
                chart_found = True
                chart_index += 1
        
        if not chart_found:
            print("  - Tidak ada chart di slide ini")

    total_charts = sum(1 for slide in presentation.slides for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.CHART)
    print(f"\nTotal chart di seluruh presentasi: {total_charts}")


def generate_chart(presentation, data, s_idx, c_idx):
    data_chart = data
    if not data_chart:
        categories = []
        values = []
    else:
        categories = []
        values = []

    sorted_data = sorted(data_chart, key=lambda x: x['key'], reverse=False)
    for item in sorted_data:
        dt = datetime.fromtimestamp(item['key'] / 1000).strftime("%d-%m-%Y")
        categories.append(dt)
        values.append(item['doc_count'])
    
    chart_data = CategoryChartData()
    if categories:
        chart_data.categories = categories
        chart_data.add_series("", values)
    else:
        chart_data.categories = [""]
        chart_data.add_series("", [0])
    
    slide_index = s_idx
    chart_index_to_replace = c_idx
    
    for i, slide in enumerate(presentation.slides):
        if i == slide_index:
            replace_chart(slide, chart_index_to_replace, chart_data)
            break

def convert_epoch_to_iso(epoch):
    # print("epoch:", epoch)
    return datetime.utcfromtimestamp(epoch / 1000).strftime('%Y-%m-%dT%H:%M:%S.%fZ')

def rank_sum(data_count, idx):
    data_rank = 0
    for avg in data_count:
        data_rank += avg
        devide = data_rank / len(data_count)
    data = data_count[idx] - devide
    return f"{round(data, 2)}"

def update_query_with_timeframe(query, params: dict):

    from_time = convert_epoch_to_iso(params["timeframe"].get("from"))
    to_time = convert_epoch_to_iso(params["timeframe"].get("to"))
    timeframe_field = params["timeframe"].get("timeframe_field")

    existing_filters = query['query']['bool']['filter']
    existing_filters = [f for f in existing_filters if not ('range' in f and timeframe_field in f['range'])]

    range_filter = {
        "range": {
            "created_at": {
                "format": "strict_date_optional_time",
                "gte": from_time,
                "lte": to_time
            }
        }
    }
    existing_filters.append(range_filter)
    query['query']['bool']['filter'] = existing_filters

    return query

def update_query_with_timeframe_must(query, params: dict):
    print("param timeframe: ", params)
    # params = parse_report_params(data)
    from_time = convert_epoch_to_iso(params["timeframe"].get("from"))
    to_time = convert_epoch_to_iso(params["timeframe"].get("to"))
    timeframe_field = params["timeframe"].get("timeframe_field")

    existing_filters = query['query']['bool']['must']
    existing_filters = [f for f in existing_filters if not ('range' in f and timeframe_field in f['range'])]

    range_filter = {
        "range": {
            "created_at": {
                "format": "strict_date_optional_time",
                "gte": from_time,
                "lte": to_time
            }
        }
    }

    print("range filter: ", range_filter)

    existing_filters.append(range_filter)
    query['query']['bool']['must'] = existing_filters

    return query

def handle_query_result_kpi(result):
    return [{"key": row[0], "dpi": row[1], "dii": row[2], "ppi": row[3], "kpi": row[4]} for row in result]

def handle_query_result_dpi(result):
    return [{"key": row[0], "exsposure": row[1], "engagement": row[2], "dpi": row[3]} for row in result]

def handle_query_result_dii(result):
    return [{"key": row[0], "exsposure": row[1], "engagement": row[2], "sentiment_exsposure": [3], "sentiment_engagement": [4], "dii": [5]} for row in result]

def handle_query_result_ppi(result):
    return [{"key": row[0], "exsposure": row[1], "engagement": row[2], "sentiment_exsposure": [3], "sentiment_engagement": [4], "ppi": [5]} for row in result]

def generate_chart_string_key(presentation, data, s_idx, c_idx):
    data_chart = data
    if not data_chart:
        categories = []
        values = []
    else:
        categories = []
        values = []

    sorted_data = sorted(data_chart, key=lambda x: x['doc_count'], reverse=False)
    for item in sorted_data:
        categories.append(item['key'])
        values.append(item['doc_count'])
    
    chart_data = CategoryChartData()
    if categories:
        chart_data.categories = categories
        chart_data.add_series("", values)
    else:
        chart_data.categories = [""]
        chart_data.add_series("", [0])
    
    slide_index = s_idx
    chart_index_to_replace = c_idx
    
    for i, slide in enumerate(presentation.slides):
        if i == slide_index:
            replace_chart(slide, chart_index_to_replace, chart_data)
            break

def process_elastic_data(elastic_data):
    rows = []
    
    # Extract the main "satwil" aggregation
    for satwil_bucket in elastic_data['aggregations']['satwil']['buckets']:
        satwil_name = satwil_bucket['key']
        
        # Extract the campaigns within the satwil bucket
        for campaign_bucket in satwil_bucket['campaign']['buckets']:
            campaign_name = campaign_bucket['key']
            
            # Extract platforms within the campaign
            for platform_bucket in campaign_bucket['platform']['buckets']:
                platform_name = platform_bucket['key']
                
                # Extract tasks within the platform
                for task_bucket in platform_bucket['task']['buckets']:
                    task_name = task_bucket['key']
                    
                    # Extract the created_at date
                    for date_bucket in task_bucket['created_at']['buckets']:
                        date_str = date_bucket['key_as_string']
                        rows.append({
                            'Satwil': satwil_name,
                            'Campaign': campaign_name,
                            'Platform': platform_name,
                            'Task': task_name,
                            'Date': date_str,
                            'Doc Count': date_bucket['doc_count']
                        })
    
    # Convert the rows into a DataFrame
    df = pd.DataFrame(rows)
    return df

def process_postgre_data(result):
    df = pd.DataFrame(result)
    return df